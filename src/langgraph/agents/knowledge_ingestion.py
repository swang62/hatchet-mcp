"""Knowledge ingestion LangGraph graph: extract → deep_inspect → chunk → store."""

import json
import os
import shutil
from datetime import datetime, timezone
from pathlib import Path
from typing import Literal, TypedDict

import numpy as np
from chromadb import PersistentClient
from chromadb.api.types import Metadata
from langchain_openai import ChatOpenAI
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_voyageai import VoyageAIEmbeddings
from langgraph.graph import END, START, StateGraph
from pydantic import SecretStr

from src.shared.constants import (
    CHROMA_DIR,
    CHUNK_OVERLAP,
    CHUNK_SIZE,
    COLLECTION_NAME,
    FILE_DIR,
    HTML_EXTENSIONS,
    IMAGE_EXTENSIONS,
    INDEX_PATH,
    OFFICE_EXTENSIONS,
    TEXT_EXTENSIONS,
    VOYAGE_MODEL,
)


class IngestionState(TypedDict):
    file_path: str
    document_id: str
    source: str
    file_type: str
    text: str
    title: str
    summary: str
    sections: list[dict]
    keywords: list[str]
    topic: str
    doc_type: str
    entities: list[str]
    chunks: list[str]
    num_chunks: int
    # Robustness: error flag and retry counters
    error: str
    inspect_retries: int
    max_retries: int


def extract_text(state: IngestionState) -> dict:
    path = Path(state["file_path"])
    ext = path.suffix.lower()
    try:
        text: str = ""
        title: str = ""

        if ext == ".pdf":
            import fitz

            doc = fitz.open(state["file_path"])
            text = "".join(str(page.get_text()) for page in doc)
            title = ((doc.metadata or {}).get("title") or "").strip()
            if not title:
                title = path.stem
            doc.close()

        elif ext in OFFICE_EXTENSIONS:
            if ext == ".docx":
                from docx import Document

                doc = Document(state["file_path"])
                text = "\n".join(p.text for p in doc.paragraphs)
            elif ext == ".pptx":
                from pptx import Presentation

                prs = Presentation(state["file_path"])
                text = "\n".join(
                    shape.text  # type: ignore[union-attr]
                    for slide in prs.slides
                    for shape in slide.shapes
                    if hasattr(shape, "text")
                )
            elif ext == ".xlsx":
                from openpyxl import load_workbook

                wb = load_workbook(state["file_path"], read_only=True, data_only=True)
                rows = []
                for sheet in wb.sheetnames:
                    for row in wb[sheet].iter_rows(values_only=True):
                        cells = [str(c) for c in row if c is not None]
                        if cells:
                            rows.append(" | ".join(cells))
                text = "\n".join(rows)
                wb.close()
            title = path.stem

        elif ext in HTML_EXTENSIONS:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="replace"), "lxml")
            for tag in soup(["script", "style", "nav", "header", "footer", "aside"]):
                tag.decompose()
            text = soup.get_text(separator="\n")
            title = soup.title.string.strip() if soup.title and soup.title.string else path.stem

        elif ext in IMAGE_EXTENSIONS:
            import easyocr

            reader = easyocr.Reader(["en"], gpu=False)
            result = reader.readtext(str(path))
            text = "\n".join(item[1] for item in result) if result else ""  # type: ignore[misc]
            title = path.stem

        elif ext in TEXT_EXTENSIONS:
            text = path.read_text(encoding="utf-8", errors="replace")
            title = path.stem

        else:
            raise ValueError(f"Unsupported file type: {ext}")

        if not text.strip():
            raise ValueError("Extracted text is empty")

        return {"text": text, "title": title, "error": ""}
    except Exception as e:
        return {"error": f"extract_text failed: {e}"}


def _llm() -> ChatOpenAI:
    return ChatOpenAI(
        model=os.getenv("OPENAI_MODEL", "gpt-4o-mini"),
        api_key=SecretStr(os.environ["OPENAI_API_KEY"]) if "OPENAI_API_KEY" in os.environ else None,
        base_url=os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1"),
        temperature=0,
    )


def deep_inspect(state: IngestionState) -> dict:
    try:
        sample = state["text"][:8000]
        llm = _llm()
        prompt = (
            "You are a research document analyst. Analyze this document and extract:\n"
            "1. A 2-3 sentence summary\n"
            "2. The main topic (short label, 2-5 words)\n"
            "3. Document type (paper/tutorial/report/manual/blog/data/other)\n"
            "4. 5-10 key technical keywords\n"
            "5. Key entities mentioned (methods, algorithms, datasets, models, people)\n"
            "6. Section headings found in the text (in order)\n\n"
            f"Document excerpt:\n{sample}\n\n"
            "Reply as JSON:\n"
            '{\n  "summary": "...",\n  "topic": "...",\n  "doc_type": "...",\n'
            '  "keywords": ["...", "..."],\n  "entities": ["...", "..."],\n'
            '  "sections": [{"heading": "...", "level": 1}]\n}'
        )
        rsp = llm.invoke(prompt)
        content = rsp.content if isinstance(rsp.content, str) else str(rsp.content)
        data = json.loads(content)

        result = {
            "summary": data.get("summary", ""),
            "topic": data.get("topic", ""),
            "doc_type": data.get("doc_type", "unknown"),
            "keywords": data.get("keywords", []),
            "entities": data.get("entities", []),
            "sections": data.get("sections", []),
            "inspect_retries": state.get("inspect_retries", 0) + 1,
            "error": "",
        }

        if not result["summary"]:
            result["error"] = "deep_inspect returned empty summary"

        return result
    except Exception as e:
        return {
            "inspect_retries": state.get("inspect_retries", 0) + 1,
            "error": f"deep_inspect failed: {e}",
        }


def chunk_text(state: IngestionState) -> dict:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    chunks = splitter.split_text(state["text"])
    return {"chunks": chunks, "num_chunks": len(chunks)}


def store_in_chroma(state: IngestionState) -> dict:
    client = PersistentClient(str(CHROMA_DIR))
    collection = client.get_or_create_collection(COLLECTION_NAME)

    file_name = Path(state["file_path"]).name
    dest = FILE_DIR / file_name
    FILE_DIR.mkdir(parents=True, exist_ok=True)
    shutil.copy2(state["file_path"], str(dest))

    model = VoyageAIEmbeddings(model=VOYAGE_MODEL)
    raw_embeddings = model.embed_documents(state["chunks"])
    embeddings = [np.array(e, dtype=np.float32) for e in raw_embeddings]

    ids = [f"{state['document_id']}_{i}" for i in range(state["num_chunks"])]
    metadatas: list[Metadata] = [
        {
            "document_id": state["document_id"],
            "source": state["source"],
            "file_type": state["file_type"],
            "chunk_index": str(i),
            "topic": state["topic"],
            "doc_type": state["doc_type"],
            "file_name": file_name,
            "file_path": str(dest),
            "title": state.get("title", ""),
            "keywords": ", ".join(state.get("keywords", [])),
            "entities": ", ".join(state.get("entities", [])),
        }
        for i in range(state["num_chunks"])
    ]

    collection.add(
        ids=ids,
        documents=state["chunks"],
        embeddings=embeddings,  # type: ignore[arg-type]
        metadatas=metadatas,
    )

    entry = {
        "document_id": state["document_id"],
        "file_name": file_name,
        "file_type": state["file_type"],
        "file_path": str(dest),
        "ingested_at": datetime.now(timezone.utc).isoformat(),
        "file_size": dest.stat().st_size,
        "num_chunks": state["num_chunks"],
        "title": state.get("title", ""),
        "summary": state["summary"],
        "topic": state["topic"],
        "keywords": state["keywords"],
        "doc_type": state["doc_type"],
        "entities": state["entities"],
        "sections": state["sections"],
    }

    INDEX_PATH.parent.mkdir(parents=True, exist_ok=True)
    index = json.loads(INDEX_PATH.read_text()) if INDEX_PATH.exists() else {}
    index[state["document_id"]] = entry
    INDEX_PATH.write_text(json.dumps(index, indent=2))

    return {"result": f"Stored {state['num_chunks']} chunks for {state['document_id']}"}


def route_after_extract(state: IngestionState) -> Literal["deep_inspect", "abort"]:
    """If extraction failed, abort. Otherwise proceed to deep_inspect."""
    if state.get("error"):
        return "abort"
    return "deep_inspect"


def route_after_inspect(state: IngestionState) -> Literal["chunk_text", "deep_inspect", "abort"]:
    """If inspection failed and retries remain, retry. Otherwise abort or proceed."""
    if state.get("error"):
        if state.get("inspect_retries", 0) < state.get("max_retries", 2):
            return "deep_inspect"
        return "abort"
    return "chunk_text"


def abort(state: IngestionState) -> dict:
    """Log the failure and terminate."""
    error = state.get("error", "Unknown error")
    print(f"[abort] Ingestion failed for {state.get('document_id', 'unknown')}: {error}")
    return {"error": error}


graph = (
    StateGraph(IngestionState)  # type: ignore[invalid-argument-type]
    .add_node("extract_text", extract_text)
    .add_node("deep_inspect", deep_inspect)
    .add_node("chunk_text", chunk_text)
    .add_node("store_in_chroma", store_in_chroma)
    .add_node("abort", abort)
    .add_edge(START, "extract_text")
    .add_conditional_edges(
        "extract_text",
        route_after_extract,
        {
            "deep_inspect": "deep_inspect",
            "abort": "abort",
        },
    )
    .add_conditional_edges(
        "deep_inspect",
        route_after_inspect,
        {
            "chunk_text": "chunk_text",
            "deep_inspect": "deep_inspect",
            "abort": "abort",
        },
    )
    .add_edge("chunk_text", "store_in_chroma")
    .add_edge("store_in_chroma", END)
    .add_edge("abort", END)
    .compile()
)
